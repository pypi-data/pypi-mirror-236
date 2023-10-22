from pathlib import Path
 
class GenericDataLoader:
    @staticmethod        
    def read_file(file_path:Path):
        if file_path.suffix ==".pdf":
            return GenericDataLoader.read_pdf_file(file_path)
        elif file_path.suffix == ".docx":
            return GenericDataLoader.read_docx_file(file_path)
        elif file_path.suffix == ".json":
            return GenericDataLoader.read_json_file(file_path)
        elif file_path.suffix == ".html":
            return GenericDataLoader.read_html_file(file_path)
        
        elif file_path.suffix == ".pptx":
            return GenericDataLoader.read_pptx_file(file_path)
        if file_path.suffix in [".txt", ".rtf", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat"]:
            return GenericDataLoader.read_text_file(file_path)
        elif file_path.suffix == ".csv":
            return GenericDataLoader.read_csv_file(file_path)
        elif file_path.suffix == ".xlsx":
            return GenericDataLoader.read_excel_file(file_path)
        else:
            raise ValueError("Unknown file type")
    def get_supported_file_types():
        return ["pdf", "txt", "docx", "json", "html", "pptx",".txt", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat", ".rtf"]    
    @staticmethod        
    def read_pdf_file(file_path):
        import PyPDF2
        def extract_text_from_pdf(file_path):
            text = ""
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            return text
        # Extract text from the PDF
        text = extract_text_from_pdf(file_path)

        # Convert to Markdown (You may need to implement custom logic based on your specific use case)
        markdown_text = text.replace('\n', '  \n')  # Adding double spaces at the end of each line for Markdown line breaks
        
        return markdown_text

    @staticmethod
    def read_docx_file(file_path):
        from docx import Document

        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        import json
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        import pandas as pd
        try:
            df = pd.read_csv(file_path)
            return df
        except FileNotFoundError:
            print(f"File '{file_path}' not found.")
            return None

    @staticmethod
    def read_excel_file(file_path, sheet_name=None):
        import pandas as pd
        try:
            df = pd.read_excel(file_path, sheet_name)
            return df
        except FileNotFoundError:
            print(f"File '{file_path}' not found.")
            return None

    @staticmethod
    def read_html_file(file_path):
        from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    
    @staticmethod
    def read_pptx_file(file_path):
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content
